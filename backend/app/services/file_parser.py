import io
from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document

async def extract_text_from_file(file):
    content = await file.read()

    if file.filename.endswith('.pdf'):
        return parse_pdf(content)
    elif file.filename.endswith('.txt'):
        return content.decode('utf-8')
    elif file.filename.endswith('.docx'):
        return parse_docx(content)
    elif file.filename.endswith('.pptx'):
        return parse_pptx(content)
    else:
        raise ValueError("Unsupported file type")
    
def parse_pdf(binary_data):
    pdf_reader = PdfReader(io.BytesIO(binary_data))
    text = ""
    for page in pdf_reader.pages:
        text += page.extract_text() + "\n"
    return text

def parse_docx(binary_data):
    document = Document(io.BytesIO(binary_data))
    text = ""
    for para in document.paragraphs:
        text += para.text + "\n"
    return text

def parse_pptx(binary_data):
    presentation = Presentation(io.BytesIO(binary_data))
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text


